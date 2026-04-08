"""
file_processor.py  (security-hardened)
---------------------------------------
Security fixes:
  [CRITICAL] Path traversal  — UUID-prefixed filenames, stem sanitized, realpath check
  [CRITICAL] File size limit — MAX_FILE_BYTES enforced BEFORE disk write
  [MEDIUM]   Filename collision — uuid4 hex prefix on every saved file
  [LOW]      Magic-byte validation — rejects disguised/polyglot files
  [LOW]      Chunk + page caps — prevents memory exhaustion on huge files
"""

import os
import re
import uuid
from pathlib import Path

UPLOAD_DIR          = "/tmp/revai_uploads"
MAX_FILE_BYTES      = 20 * 1024 * 1024   # 20 MB
MAX_CHUNKS_PER_FILE = 500
MAX_PDF_PAGES       = 200
MAX_OCR_PAGES       = 20
MAX_SLIDES          = 100
MAX_SHEETS          = 20
MAX_ROWS_PER_SHEET  = 5000
MAX_TEXT_READ       = 5 * 1024 * 1024    # 5 MB for plain text

os.makedirs(UPLOAD_DIR, exist_ok=True)

SUPPORTED  = {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt", ".csv", ".md", ".xlsx"}
CHUNK_SIZE = 1000
OVERLAP    = 100


# ─────────────────────────────────────────────
# Internal helpers
# ─────────────────────────────────────────────

def _safe_filename(original: str) -> str:
    """
    Strip path components, sanitize stem, prepend UUID.
    '../../etc/passwd.pdf' → 'a3f1c2d4_etc_passwd.pdf'
    """
    name      = Path(original).name           # drops any directory part
    ext       = Path(name).suffix.lower()
    stem      = Path(name).stem
    safe_stem = re.sub(r"[^a-zA-Z0-9_\-]", "_", stem)[:40] or "file"
    return f"{uuid.uuid4().hex[:8]}_{safe_stem}{ext}"


def _check_magic(file_bytes: bytes, ext: str) -> bool:
    """
    Compare file header bytes against expected magic bytes.
    Prevents polyglot / disguised executables.
    """
    h = file_bytes[:8]

    if ext == ".pdf":
        return h[:4] == b"%PDF"

    if ext in (".pptx", ".ppt", ".docx", ".xlsx"):
        # All modern Office formats are ZIP archives
        return h[:4] == b"PK\x03\x04"

    if ext == ".doc":
        # Legacy OLE2 Compound Document
        return h[:4] == b"\xd0\xcf\x11\xe0"

    if ext in (".txt", ".md", ".csv"):
        # Reject obvious binary files (null bytes in first 512)
        return b"\x00" not in file_bytes[:512]

    return True


# ─────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────

def save_upload(file_bytes: bytes, filename: str) -> str:
    """
    Validate and save uploaded bytes. Returns safe disk path.
    Raises ValueError for all invalid inputs (caller should catch and return 400).
    """
    # 1 — Extension whitelist
    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Allowed: {', '.join(sorted(SUPPORTED))}"
        )

    # 2 — Empty file
    if not file_bytes:
        raise ValueError("Uploaded file is empty.")

    # 3 — Size limit (checked BEFORE writing to disk)
    if len(file_bytes) > MAX_FILE_BYTES:
        mb = len(file_bytes) / (1024 * 1024)
        limit_mb = MAX_FILE_BYTES // (1024 * 1024)
        raise ValueError(f"File too large ({mb:.1f} MB). Limit is {limit_mb} MB.")

    # 4 — Magic byte check (detects disguised files)
    if not _check_magic(file_bytes, ext):
        raise ValueError(
            f"File content does not match declared extension '{ext}'. "
            "The file may be corrupted or intentionally disguised."
        )

    # 5 — Safe filename (UUID prefix + sanitized stem)
    safe_name = _safe_filename(filename)
    path      = os.path.join(UPLOAD_DIR, safe_name)

    # 6 — Belt-and-suspenders: confirm final path stays inside UPLOAD_DIR
    if not os.path.realpath(path).startswith(os.path.realpath(UPLOAD_DIR)):
        raise ValueError("Invalid file path — path traversal detected.")

    with open(path, "wb") as f:
        f.write(file_bytes)

    print(f"[file_processor] Saved: {path} ({len(file_bytes):,} bytes)")
    return path


def extract_text_chunks(file_path: str) -> list:
    """Parse a file and return a capped list of text chunk strings."""
    ext = Path(file_path).suffix.lower()
    print(f"[file_processor] Parsing '{Path(file_path).name}' (ext={ext})")

    parsers = {
        ".pdf":  _parse_pdf,
        ".pptx": _parse_pptx, ".ppt": _parse_pptx,
        ".docx": _parse_docx, ".doc": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".txt":  _parse_text, ".md": _parse_text, ".csv": _parse_text,
    }
    parser = parsers.get(ext)
    if not parser:
        raise ValueError(f"No parser for extension: {ext}")

    raw    = parser(file_path)
    chunks = _split_into_chunks(raw, CHUNK_SIZE, OVERLAP)

    if len(chunks) > MAX_CHUNKS_PER_FILE:
        print(f"[file_processor] Capping {len(chunks)} → {MAX_CHUNKS_PER_FILE} chunks")
        chunks = chunks[:MAX_CHUNKS_PER_FILE]

    print(f"[file_processor] Done — {len(chunks)} chunks.")
    return chunks


# ─────────────────────────────────────────────
# Parsers
# ─────────────────────────────────────────────

def _parse_pdf(path: str) -> str:
    import pdfplumber
    pages_text = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages[:MAX_PDF_PAGES]):
            text = page.extract_text()
            if text and text.strip():
                pages_text.append(f"[Page {i+1}]\n{text.strip()}")
    if pages_text:
        return "\n\n".join(pages_text)
    return _ocr_pdf(path)


def _ocr_pdf(path: str) -> str:
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError as e:
        raise ImportError(f"OCR dependencies missing: {e}") from e

    images    = convert_from_path(path, dpi=150, last_page=MAX_OCR_PAGES)
    ocr_pages = []
    for i, img in enumerate(images):
        text = pytesseract.image_to_string(img, lang="eng").strip()
        if text:
            ocr_pages.append(f"[Page {i+1} — OCR]\n{text}")
    if not ocr_pages:
        raise ValueError("OCR produced no text. PDF may be corrupted or image-only.")
    return "\n\n".join(ocr_pages)


def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs    = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides[:MAX_SLIDES], 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    if not slides:
        raise ValueError("No text found in PowerPoint file.")
    return "\n\n".join(slides)


def _parse_docx(path: str) -> str:
    from docx import Document
    doc   = Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)
    if not parts:
        raise ValueError("No text found in Word document.")
    return "\n\n".join(parts)


def _parse_xlsx(path: str) -> str:
    import openpyxl
    wb   = openpyxl.load_workbook(path, data_only=True, read_only=True)
    rows = []
    for sheet in list(wb.worksheets)[:MAX_SHEETS]:
        rows.append(f"[Sheet: {sheet.title}]")
        row_count = 0
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                rows.append(row_text)
                row_count += 1
                if row_count >= MAX_ROWS_PER_SHEET:
                    rows.append("[...truncated — row limit reached...]")
                    break
    wb.close()
    if not rows:
        raise ValueError("No data found in Excel file.")
    return "\n".join(rows)


def _parse_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read(MAX_TEXT_READ)
    if not text.strip():
        raise ValueError("File appears to be empty.")
    return text


# ─────────────────────────────────────────────
# Chunker
# ─────────────────────────────────────────────

def _split_into_chunks(text: str, size: int, overlap: int) -> list:
    chunks, start = [], 0
    while start < len(text):
        chunk = text[start : start + size].strip()
        if chunk:
            chunks.append(chunk)
        start += size - overlap
    return chunks